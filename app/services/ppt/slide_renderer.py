# app\services\ppt\slide_renderer.py
import os
import json
from pptx import Presentation
from .text_renderer import render_text
from .chart_renderer import render_chart
from .image_renderer import render_image

PROJECT_ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), "../../.."))
TEMPLATE_DIR = os.path.join(PROJECT_ROOT, "templates_library", "ppt_templates")
TEMPLATE_SPEC_DIR = os.path.join(PROJECT_ROOT, "templates_library", "template_specs")

# Load master template specs once
MASTER_TEMPLATE_JSON_PATH = os.path.join(TEMPLATE_SPEC_DIR, "master_template.json")
with open(MASTER_TEMPLATE_JSON_PATH, "r", encoding="utf-8") as f:
    MASTER_TEMPLATE_SPEC = json.load(f)


# Build mapping: (layout_name, placeholder_name) -> placeholder_index
PLACEHOLDER_ID_MAP = {}
for layout in MASTER_TEMPLATE_SPEC["layouts"]:
    layout_name = layout["layout_name"]
    for placeholder_name, ph_info in layout["placeholders"].items():
        PLACEHOLDER_ID_MAP[(layout_name, placeholder_name)] = ph_info["placeholder_index"]


def get_template(template_id):
    template_path = os.path.join(TEMPLATE_DIR, f"{template_id}.pptx")
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template file not found: {template_path}")
    return Presentation(template_path)


def render_slide(prs, slide_json):
    layout_name = slide_json["template_id"]
    content = slide_json["content"]

    # Find layout in master template by name
    layout = next(l for l in prs.slide_layouts if l.name == layout_name)
    slide = prs.slides.add_slide(layout)

    # Map content keys to placeholder indexes
    for placeholder_name, value in content.items():
        placeholder_index = PLACEHOLDER_ID_MAP.get((layout_name, placeholder_name))
        if placeholder_index is None:
            print(f"[WARN] Placeholder '{placeholder_name}' not found in layout '{layout_name}'")
            continue

        # Find placeholder in slide by idx
        shape = next((s for s in slide.placeholders if s.placeholder_format.idx == placeholder_index), None)
        if shape is None:
            print(f"[WARN] Placeholder idx '{placeholder_index}' not found in slide '{layout_name}'")
            continue

        # Render content based on type
        if hasattr(shape, "text_frame"):
            render_text(shape, value)
        elif shape.placeholder_format.type.name == "PICTURE":
            render_image(shape, value)
        elif shape.placeholder_format.type.name == "CHART":
            render_chart(shape, value)
        # Add other types as needed
