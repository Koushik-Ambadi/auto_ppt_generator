# Create 5 separately downloadable, well-designed PPT template files
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL

def add_title_slide(filename):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    # Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(18, 32, 47)
    
    # Title
    title = slide.shapes.title
    title.text = "TITLE PLACEHOLDER"
    title_tf = title.text_frame.paragraphs[0]
    title_tf.font.size = Pt(44)
    title_tf.font.bold = True
    title_tf.font.color.rgb = RGBColor(255, 255, 255)
    
    # Subtitle Box
    box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1.5))
    tf = box.text_frame
    tf.text = "Subtitle Placeholder"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
    
    prs.save(f"templates_library\\ppt_templates\\{filename}")

def add_exec_summary(filename):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    slide.shapes.title.text = "Executive Summary"
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide.placeholders[1]
    content.text = "• Key Point One\n• Key Point Two\n• Key Point Three\n\nSummary paragraph placeholder."
    for p in content.text_frame.paragraphs:
        p.font.size = Pt(20)
    
    prs.save(f"templates_library\\ppt_templates\\{filename}")

def add_two_column(filename):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    slide.shapes.title.text = "Two Column Layout"
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(34)
    
    # Left text
    left = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4.2), Inches(4))
    left_tf = left.text_frame
    left_tf.text = "• Bullet One\n• Bullet Two\n• Bullet Three"
    for p in left_tf.paragraphs:
        p.font.size = Pt(20)
    
    # Right placeholder box
    right = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(5.2), Inches(1.8), Inches(4.2), Inches(4)
    )
    right.fill.background()
    right.line.color.rgb = RGBColor(120, 120, 120)
    right.text_frame.text = "Image / Graphic Placeholder"
    
    prs.save(f"templates_library\\ppt_templates\\{filename}")

def add_chart_slide(filename):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    slide.shapes.title.text = "Chart / Data Slide"
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(34)
    
    chart_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(1), Inches(1.8), Inches(8), Inches(4)
    )
    chart_box.fill.background()
    chart_box.line.color.rgb = RGBColor(150, 150, 150)
    chart_box.text_frame.text = "Chart Placeholder"
    
    prs.save(f"templates_library\\ppt_templates\\{filename}")

def add_image_slide(filename):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    slide.shapes.title.text = "Image Focus Slide"
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(34)
    
    image_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(1), Inches(1.5), Inches(8), Inches(4.5)
    )
    image_box.fill.background()
    image_box.line.color.rgb = RGBColor(100, 100, 100)
    image_box.text_frame.text = "Full Width Image Placeholder"
    
    prs.save(f"templates_library\\ppt_templates\\{filename}")

# Generate files
add_title_slide("TEMPLATE_TITLE_V1.pptx")
add_exec_summary("TEMPLATE_EXEC_SUMMARY_V1.pptx")
add_two_column("TEMPLATE_TWO_COLUMN_V1.pptx")
add_chart_slide("TEMPLATE_CHART_V1.pptx")
add_image_slide("TEMPLATE_IMAGE_V1.pptx")

print("All 5 individual templates created successfully.")
