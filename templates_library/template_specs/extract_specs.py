# templates_library\template_specs\extract_specs.py

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
import json
import os

# -----------------------------
# Placeholder → semantic field map
# -----------------------------
PLACEHOLDER_FIELD_MAP = {
    PP_PLACEHOLDER.TITLE: "title",
    PP_PLACEHOLDER.CENTER_TITLE: "title",
    PP_PLACEHOLDER.SUBTITLE: "subtitle",
    PP_PLACEHOLDER.BODY: "text",        # generic text container
    PP_PLACEHOLDER.OBJECT: "multi_media",
    PP_PLACEHOLDER.PICTURE: "image",
    PP_PLACEHOLDER.CHART: "chart",
    PP_PLACEHOLDER.TABLE: "table",
    PP_PLACEHOLDER.FOOTER: "footnote"
}

# -----------------------------
# Utility functions
# -----------------------------
def infer_field_type(ph):
    """Returns the field type and extra info for a placeholder"""
    ph_type = ph.placeholder_format.type

    if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.SUBTITLE, PP_PLACEHOLDER.FOOTER):
        return "text", {}
    if ph_type == PP_PLACEHOLDER.BODY:
        return "text", {"supports_bullets": True, "supports_paragraphs": True}
    if ph_type == PP_PLACEHOLDER.PICTURE:
        return "image", {}
    if ph_type == PP_PLACEHOLDER.CHART:
        return "chart", {}
    if ph_type == PP_PLACEHOLDER.TABLE:
        return "table", {}
    return "text", {}

# -----------------------------
# Extract master template schema
# -----------------------------
def extract_master_template_schema(ppt_path):
    prs = Presentation(ppt_path)
    template_id = os.path.splitext(os.path.basename(ppt_path))[0]

    schema = {
        "template_id": template_id,
        "layouts": []
    }

    for layout_index, layout in enumerate(prs.slide_layouts):
        layout_schema = {
            "layout_index": layout_index,
            "layout_name": layout.name,
            "placeholders": {}
        }

        for ph in layout.placeholders:
            ph_index = ph.placeholder_format.idx
            ph_type = ph.placeholder_format.type
            placeholder_name = ph.name

            # Ensure unique placeholder names in layout
            original_name = placeholder_name
            counter = 2
            while placeholder_name in layout_schema["placeholders"]:
                placeholder_name = f"{original_name}_{counter}"
                counter += 1

            field_type, extra = infer_field_type(ph)

            layout_schema["placeholders"][placeholder_name] = {
                "placeholder_name": placeholder_name,       # editable name
                "placeholder_index": ph_index,             # PPT ID
                "placeholder_type": ph_type.name if hasattr(ph_type, "name") else str(ph_type),
                "field_type": field_type,
                **extra,
                "left": ph.left,
                "top": ph.top,
                "width": ph.width,
                "height": ph.height
            }

        schema["layouts"].append(layout_schema)

    return schema

# -----------------------------
# Main runner
# -----------------------------
if __name__ == "__main__":
    INPUT_PATH = "templates_library/ppt_templates/master_template.pptx"
    OUTPUT_DIR = "templates_library/template_specs"
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    if not os.path.isfile(INPUT_PATH):
        raise FileNotFoundError(f"Master template not found: {INPUT_PATH}")

    schema = extract_master_template_schema(INPUT_PATH)

    output_path = os.path.join(OUTPUT_DIR, f"{schema['template_id']}.json")
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(schema, f, indent=2)

    print(f"Master template schema generated → {output_path}")
